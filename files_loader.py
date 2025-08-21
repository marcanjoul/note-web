import os  # Provides functions for interacting with the operating system, such as file path manipulation.
import fitz  # PyMuPDF, used to read and process PDF files.
from docx import Document  # Used to read and process Microsoft Word (.docx) files.
from openpyxl import load_workbook  # Used to read and process Microsoft Excel (.xlsx) files.
from pptx import Presentation  # Used to read and process Microsoft PowerPoint (.pptx) files.

def extract_text_from_pdf(file_path):
    """
    Opens a PDF file and extracts all the text from every page.

    Args:
        file_path (str): The path to the PDF file.

    Returns:
        str: The full extracted text from the PDF.
    """
    text = ""  # Initialize an empty string to collect text from all pages.

    try:
        # Open the PDF file using PyMuPDF's fitz.open, which returns a document object.
        with fitz.open(file_path) as doc:
            # Iterate through each page in the PDF document.
            for page in doc:
                # Extract text from the current page and append it to the running string.
                text += page.get_text()
    except Exception as e:
        # Print an error message if something goes wrong during reading.
        print(f" Error reading {file_path}: {e}")

    # Return the accumulated text from all pages.
    return text

def extract_text_from_multiple_pdfs(file_paths):
    """
    Extracts and combines text from multiple PDF files.

    Args:
        file_paths (List[str]): A list of paths to PDF files.

    Returns:
        str: Combined text from all PDFs.
    """
    combined_text = ""  # Initialize an empty string to collect text from all PDFs.
    for path in file_paths:
        # Print the path of the PDF currently being processed.
        print(f" Extracting: {path}")
        # Extract text from the current PDF and append it to the combined string, separated by a newline.
        combined_text += extract_text_from_pdf(path) + "\n"
    # Return the combined text from all PDFs.
    return combined_text

def extract_text_from_docx(file_path):
    """
    Opens a Word document (.docx) and extracts all the text from every paragraph.

    Args:
        file_path (str): The path to the docx file.

    Returns:
        str: The full extracted text from the docx.
    """
    text = ""  # Initialize an empty string to collect text from all paragraphs.

    try:
        # Load the Word document using python-docx's Document class.
        doc = Document(file_path)
        # Iterate through each paragraph in the document.
        for para in doc.paragraphs:
            # Append the paragraph text and a newline to the running string.
            text += para.text + "\n"
    except Exception as e:
        # Print an error message if something goes wrong during reading.
        print(f"Error reading {file_path}: {e}")
    # Return the accumulated text from all paragraphs.
    return text

def extract_text_from_xlsx(file_path):
    """
    Opens an Excel workbook (.xlsx) and extracts all cell values as text.

    Args:
        file_path (str): The path to the xlsx file.

    Returns:
        str: The full extracted text from the xlsx.
    """
    text = ""  # Initialize an empty string to collect text from all cells.

    try:
        # Load the Excel workbook using openpyxl's load_workbook function.
        wb = load_workbook(file_path)
        # Iterate through each worksheet in the workbook.
        for sheet in wb.worksheets:
            # Iterate through each row in the worksheet, retrieving cell values only.
            for row in sheet.iter_rows(values_only=True):
                # Iterate through each cell in the row.
                for cell in row:
                    # If the cell is not empty, convert its value to a string and append it.
                    if cell:
                        text += str(cell) + " "
                # Add a newline after each row for readability.
                text += "\n"
    except Exception as e:
        # Print an error message if something goes wrong during reading.
        print(f"Error reading {file_path}: {e}")
    # Return the accumulated text from all cells.
    return text

def extract_text_from_pptx(file_path):
    """
    Opens a PowerPoint presentation (.pptx) and extracts all text from shapes on each slide.

    Args:
        file_path (str): The path to the pptx file.

    Returns:
        str: The full extracted text from the pptx.
    """
    text = ""  # Initialize an empty string to collect text from all slides.

    try:
        # Load the PowerPoint presentation using python-pptx's Presentation class.
        prs = Presentation(file_path)
        # Iterate through each slide in the presentation.
        for slide in prs.slides:
            # Iterate through each shape on the slide.
            for shape in slide.shapes:
                # If the shape has a 'text' attribute, extract and append its text.
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        # Print an error message if something goes wrong during reading.
        print(f"Error reading {file_path}: {e}")
    # Return the accumulated text from all shapes.
    return text

def extract_text_from_file(file_path):
    """
    Determines the file type by its extension and extracts text using the appropriate function.

    Args:
        file_path (str): The path to the file.

    Returns:
        str: The extracted text, or an empty string if the file type is unsupported.
    """
    # Get the file extension and convert it to lowercase for comparison.
    ext = os.path.splitext(file_path)[1].lower()

    # Call the appropriate extraction function based on the file extension.
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".xlsx":
        return extract_text_from_xlsx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        # Print a message if the file type is unsupported and return an empty string.
        print(f"Unsupported file type: {ext}")
        return ""